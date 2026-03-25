#!/usr/bin/env python3
"""
批量患者指南检索管道工具

子命令:
  parse       - 解析输入 xlsx → patients.json
  split       - 将 patients.json 分成多个批次文件
  orchestrate - 自动编排批处理流程（扫描知识库+生成 prompt）
  merge       - 合并多个 rag_batch_*.json 为 rag_results.json
  validate    - 检查 rag_results.json 质量与完整性
  generate    - 从 RAG 结果 JSON 生成 xlsx/docx/pptx 产出物
"""

from __future__ import annotations

import argparse
import json
import os
import re
import sys
from datetime import date, datetime
from pathlib import Path


# ─── parse 子命令 ─────────────────────────────────────────────────────────────


STRUCTURED_ID_COL = "患者ID号"
NARRATIVE_COL = "病情总结"

# 结构化表 26 列 → JSON 字段映射
STRUCTURED_FIELD_MAP = {
    "患者ID号": "patient_id",
    "患者姓名": "patient_name",
    "性别": "gender",
    "年龄": "age",
    "原发部位": "primary_site",
    "Siewert分型": "siewert_type",
    "病理类型": "pathology",
    "患者类型": "patient_type",
    "既往治疗说明": "prior_treatment",
    "原发灶数量": "lesion_count",
    "活检样本-分子分型": "biopsy_molecular",
    "大体样本-分子分型": "gross_molecular",
    "异常肿瘤标记物": "abnormal_markers",
    "治疗后血清肿瘤标记物变化": "marker_change",
    "分期前缀": "staging_prefix",
    "T分期": "t_stage",
    "T4b受侵脏器": "t4b_invasion",
    "N分期": "n_stage",
    "M分期": "m_stage",
    "M转移脏器": "m_sites",
    "分期备注（多病灶可在此补充说明）": "staging_notes",
    "治疗后症状变化": "symptom_change",
    "评效": "response",
    "是否合并肿瘤急症": "tumor_emergency",
    "关键合并症（如有影响诊疗决策的重大合并症，请在此处进行描述）": "comorbidities",
}


def detect_format(headers: list[str]) -> str:
    """根据列名自动检测输入格式"""
    header_set = set(h for h in headers if h)
    if STRUCTURED_ID_COL in header_set and len(headers) >= 10:
        return "structured"
    if NARRATIVE_COL in header_set and len(headers) <= 5:
        return "narrative"
    raise ValueError(
        f"无法识别输入格式。列名: {headers[:5]}...\n"
        f"支持的格式:\n"
        f"  结构化: 需要 '{STRUCTURED_ID_COL}' 列且 >=10 列\n"
        f"  半结构化: 需要 '{NARRATIVE_COL}' 列且 <=5 列"
    )


def parse_structured(ws) -> list[dict]:
    """解析结构化数据表 (26 列)"""
    rows = list(ws.iter_rows(values_only=True))
    if not rows:
        return []
    headers = [str(h) if h else "" for h in rows[0]]

    # 建立列索引
    col_idx = {}
    for i, h in enumerate(headers):
        for src, dst in STRUCTURED_FIELD_MAP.items():
            if src in h:  # 容忍列名略有不同
                col_idx[dst] = i
                break

    patients = []
    for row in rows[1:]:
        if not row or not row[col_idx.get("patient_id", 0)]:
            continue
        p = {}
        for field, idx in col_idx.items():
            val = row[idx] if idx < len(row) else None
            p[field] = str(val).strip() if val is not None else None
        # age 转整数
        if p.get("age"):
            try:
                p["age"] = int(float(p["age"]))
            except (ValueError, TypeError):
                pass
        p["clinical_narrative"] = None
        patients.append(p)
    return patients


def parse_narrative(ws) -> list[dict]:
    """解析半结构化病情总结 (3 列)"""
    rows = list(ws.iter_rows(values_only=True))
    if not rows:
        return []
    headers = [str(h) if h else "" for h in rows[0]]

    # 找到各列位置
    id_idx = next((i for i, h in enumerate(headers) if "ID" in h.upper()), 0)
    name_idx = next((i for i, h in enumerate(headers) if "姓名" in h), 1)
    narrative_idx = next((i for i, h in enumerate(headers) if "病情" in h or "总结" in h), 2)

    patients = []
    for row in rows[1:]:
        if not row or not row[id_idx]:
            continue
        p = {
            "patient_id": str(row[id_idx]).strip() if row[id_idx] else None,
            "patient_name": str(row[name_idx]).strip() if row[name_idx] else None,
            "clinical_narrative": str(row[narrative_idx]).strip() if row[narrative_idx] else None,
        }
        # 其他字段置 null，由 Claude 从 narrative 推断
        for field in STRUCTURED_FIELD_MAP.values():
            if field not in p:
                p[field] = None
        patients.append(p)
    return patients


def cmd_parse(args):
    """parse 子命令入口"""
    import openpyxl

    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"输入文件不存在: {input_path}", file=sys.stderr)
        sys.exit(1)

    wb = openpyxl.load_workbook(input_path, read_only=True, data_only=True)
    ws = wb.active

    headers = [cell.value for cell in next(ws.iter_rows(min_row=1, max_row=1))]
    fmt = detect_format(headers)
    print(f"检测到输入格式: {fmt}")

    if fmt == "structured":
        patients = parse_structured(ws)
    else:
        patients = parse_narrative(ws)

    wb.close()

    result = {
        "input_format": fmt,
        "input_file": str(input_path),
        "parsed_at": str(date.today()),
        "patient_count": len(patients),
        "patients": patients,
    }

    output_path = Path(args.output).resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"已解析 {len(patients)} 位患者 → {output_path}")


# ─── split 子命令 ─────────────────────────────────────────────────────────────


def _split_patients(patients: list[dict], batch_size: int) -> list[list[dict]]:
    """将患者列表分成多个批次（纯函数，供 split 和 orchestrate 共用）"""
    if not patients:
        return []
    return [patients[i:i + batch_size] for i in range(0, len(patients), batch_size)]


def cmd_split(args):
    """split 子命令入口 — 将 patients.json 分成多个批次文件"""
    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"输入文件不存在: {input_path}", file=sys.stderr)
        sys.exit(1)

    data = json.loads(input_path.read_text(encoding="utf-8"))
    patients = data.get("patients", [])
    batch_size = args.batch_size

    if not patients:
        print("患者列表为空，无需分批", file=sys.stderr)
        sys.exit(1)

    output_dir = Path(args.output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    batches = _split_patients(patients, batch_size)

    for idx, batch in enumerate(batches, 1):
        batch_data = {
            "input_format": data.get("input_format"),
            "input_file": data.get("input_file"),
            "parsed_at": data.get("parsed_at"),
            "batch_index": idx,
            "batch_count": len(batches),
            "patient_count": len(batch),
            "patients": batch,
        }
        batch_file = output_dir / f"batch_{idx:03d}.json"
        batch_file.write_text(
            json.dumps(batch_data, ensure_ascii=False, indent=2), encoding="utf-8"
        )

    print(f"已将 {len(patients)} 位患者分为 {len(batches)} 批（每批 {batch_size} 人）→ {output_dir}/")


# ─── orchestrate 子命令 ──────────────────────────────────────────────────────


def resolve_kb_root(explicit_path: str | None) -> Path:
    """按优先级解析知识库根路径。

    优先级: --kb-root > MEDICAL_GUIDELINES_DIR > ./guidelines/ > ./knowledge/
    验证: 目标路径下存在 data_structure.md
    """
    candidates = []

    if explicit_path:
        candidates.append(Path(explicit_path).resolve())
    else:
        env = os.environ.get("MEDICAL_GUIDELINES_DIR")
        if env:
            candidates.append(Path(env).resolve())
        candidates.append(Path("guidelines").resolve())
        candidates.append(Path("knowledge").resolve())

    for p in candidates:
        if p.is_dir() and (p / "data_structure.md").exists():
            return p

    tried = ", ".join(str(c) for c in candidates)
    print(f"无法找到知识库（需含 data_structure.md）。已尝试: {tried}", file=sys.stderr)
    sys.exit(1)


def scan_knowledge_base(kb_root: Path) -> dict:
    """自动扫描知识库结构，返回 kb_profile 字典。

    解析策略 (D1): 正则匹配表头定位 markdown 表格。
    Fallback: 解析失败时通过目录枚举发现 org。
    """
    profile = {
        "orgs": [],
        "org_files": {},
        "org_keywords": {},
        "clinical_question_map": {},
        "root_index_content": "",
    }

    root_ds = kb_root / "data_structure.md"
    root_text = ""
    if root_ds.exists():
        root_text = root_ds.read_text(encoding="utf-8")
        profile["root_index_content"] = root_text

    parsed_orgs = _parse_org_names_from_root(root_text)
    if not parsed_orgs:
        print("  ⚠ 根 data_structure.md 解析失败，fallback 到目录枚举", file=sys.stderr)
        parsed_orgs = _enumerate_org_dirs(kb_root)

    for org in parsed_orgs:
        org_dir = kb_root / org
        extracted_dir = org_dir / "extracted"
        if not extracted_dir.is_dir():
            print(f"  ⚠ {org}/ 无 extracted/ 子目录，跳过", file=sys.stderr)
            continue
        txt_files = sorted(extracted_dir.glob("*.txt"))
        if not txt_files:
            print(f"  ⚠ {org}/extracted/ 无 .txt 文件，跳过", file=sys.stderr)
            continue

        profile["orgs"].append(org)
        profile["org_files"][org] = [
            {
                "file": f.name,
                "path": str(f),
                "lines": sum(1 for _ in f.open(encoding="utf-8")),
            }
            for f in txt_files
        ]

        org_ds = org_dir / "data_structure.md"
        if org_ds.exists():
            profile["org_keywords"][org] = _parse_keywords_from_ds(
                org_ds.read_text(encoding="utf-8")
            )

    profile["clinical_question_map"] = _parse_clinical_question_map(root_text)

    print(f"  知识库扫描完成: {len(profile['orgs'])} 个组织, "
          f"{sum(len(v) for v in profile['org_files'].values())} 个文件")
    return profile


def _parse_org_names_from_root(text: str) -> list[str]:
    """从根 data_structure.md 解析组织名（### OrgName/ 模式）"""
    orgs = []
    for m in re.finditer(r'^###\s+(\w[\w-]*)/\s*$', text, re.MULTILINE):
        orgs.append(m.group(1))
    return orgs


def _enumerate_org_dirs(kb_root: Path) -> list[str]:
    """Fallback: 枚举知识库根目录下的子目录"""
    return sorted(
        d.name for d in kb_root.iterdir()
        if d.is_dir() and not d.name.startswith(".")
    )


def _parse_keywords_from_ds(text: str) -> dict[str, list[str]]:
    """从 data_structure.md 解析'常用检索关键词'区块"""
    keywords = {}
    current_category = None
    in_keyword_section = False

    for line in text.split("\n"):
        if "检索关键词" in line and line.startswith("#"):
            in_keyword_section = True
            continue
        if in_keyword_section:
            if line.startswith("---") or (line.startswith("#") and "检索" not in line):
                break
            cat_match = re.match(r'^###\s+(.+)', line)
            if cat_match:
                current_category = cat_match.group(1).strip()
                keywords[current_category] = []
                continue
            if current_category and line.startswith("- "):
                items = [k.strip() for k in line[2:].split(",")]
                keywords[current_category].extend(items)

    return keywords


def _parse_clinical_question_map(text: str) -> dict:
    """从根 data_structure.md 解析'临床问题→指南映射'表格"""
    cq_map = {}
    in_table = False

    for line in text.split("\n"):
        if "临床问题" in line and "指南映射" in line:
            in_table = True
            continue
        if in_table:
            if line.startswith("|") and "---" not in line and "临床问题" not in line:
                cols = [c.strip() for c in line.split("|")[1:-1]]
                if len(cols) >= 3:
                    question = cols[0]
                    primary = [g.strip() for g in cols[1].split(",") if g.strip()]
                    supplementary = [g.strip() for g in cols[2].split(",") if g.strip()]
                    cq_map[question] = {"primary": primary, "supplementary": supplementary}
            elif not line.startswith("|") and line.strip():
                break

    return cq_map


# grep 特殊字符转义
_GREP_SPECIAL = re.compile(r'([\[\]().*+?{}\\^$|])')


def escape_grep_keyword(keyword: str) -> list[str]:
    """转义 grep 特殊字符，含括号时额外生成去括号变体 (D4)。

    Returns: 1-2 个转义后的关键词列表
    """
    escaped = _GREP_SPECIAL.sub(r'\\\1', keyword)
    variants = [escaped]

    # D4: 含括号时生成去括号变体
    if "(" in keyword or "[" in keyword:
        stripped = re.sub(r'[()[\]]', '', keyword)
        tokens = [t for t in re.split(r'[^\w\u4e00-\u9fff]+', stripped) if t]
        if len(tokens) >= 2:
            variant = ".*".join(_GREP_SPECIAL.sub(r'\\\1', t) for t in tokens)
            variants.append(variant)

    return variants


def generate_grep_commands(
    patient_features: dict,
    kb_profile: dict,
    kb_root: "Path",
) -> list[dict]:
    """为一位患者生成覆盖所有 org 的 grep 命令。

    按临床维度合并关键词到单条 grep（用 \\| 分隔），减少命令数。
    返回: [{"org": str, "dimension": str, "command": str}, ...]
    """
    all_kw = patient_features.get("all_keywords", [])
    if not all_kw:
        return []

    dimensions = {}
    for key, val in patient_features.items():
        if key == "all_keywords" or not val:
            continue
        if key.endswith("_keywords"):
            dim_name = key.replace("_keywords", "")
            dimensions[dim_name] = val

    if not dimensions:
        dimensions["general"] = all_kw

    commands = []
    for org in kb_profile["orgs"]:
        files = kb_profile["org_files"].get(org, [])
        if not files:
            continue
        glob_pattern = str(kb_root / org / "extracted" / "*.txt")

        for dim_name, keywords in dimensions.items():
            all_variants = []
            for kw in keywords:
                all_variants.extend(escape_grep_keyword(kw))

            if not all_variants:
                continue

            pattern = "\\|".join(all_variants)
            cmd = f'grep -n -i "{pattern}" {glob_pattern}'
            commands.append({
                "org": org,
                "dimension": dim_name,
                "command": cmd,
            })

    return commands


# ─── 临床特征提取 ──────────────────────────────────────────────────────────────

_MOLECULAR_PATTERNS = re.compile(
    r'(HER2|Her2|her2|MSI-H|MSS|dMMR|pMMR|PD-L1|CPS[≥<>\d]+|EGFR|ALK|ROS1|'
    r'NTRK|BRAF|KRAS|NRAS|PIK3CA|FGFR2|Claudin[\s-]?18)',
    re.IGNORECASE,
)
_STAGING_PATTERNS = re.compile(
    r'((?:yc|c|p)?T[1-4][a-d]?|N[0-3][a-b]?|M[01]|'
    r'stage\s*(?:I{1,3}V?|IV)|[IⅠⅡⅢⅣ]+[A-C]?期)',
    re.IGNORECASE,
)
_METASTASIS_SITES = {
    "腹膜": ["腹膜", "peritoneal", "peritoneum"],
    "肝": ["肝", "liver", "hepatic"],
    "肺": ["肺", "lung", "pulmonary"],
    "骨": ["骨", "bone", "osseous"],
    "脑": ["脑", "brain", "cerebral"],
    "卵巢": ["卵巢", "ovarian", "Krukenberg"],
    "淋巴结": ["远处淋巴结", "distant lymph", "Virchow"],
}
_TREATMENT_PATTERNS = re.compile(
    r'(SOX|XELOX|CAPOX|FOLFOX|FLOT|S-1|替吉奥|卡培他滨|奥沙利铂|'
    r'PD-1|PD-L1|pembrolizumab|nivolumab|trastuzumab|'
    r'曲妥珠单抗|帕博利珠单抗|纳武利尤单抗|信迪利单抗|替雷利珠单抗|'
    r'化疗|靶向|免疫|放疗|内镜|手术|'
    r'\d+C\s+\w+)',
    re.IGNORECASE,
)
_EMERGENCY_KEYWORDS = ["出血", "bleeding", "梗阻", "obstruction", "穿孔", "perforation", "急症"]
_COMORBIDITY_PATTERNS = re.compile(
    r'(糖尿病|diabetes|高血压|hypertension|肾功能不全|renal|心[脏功]|cardiac|'
    r'肝硬化|cirrhosis|COPD|肺功能|elderly|高龄)',
    re.IGNORECASE,
)


def extract_patient_features(patient: dict) -> dict:
    """从患者数据提取 9 维临床特征关键词 (D2: 正则扫描 + confidence)。"""
    features = {
        "diagnosis_keywords": [],
        "staging_keywords": [],
        "metastasis_keywords": [],
        "molecular_keywords": [],
        "treatment_keywords": [],
        "marker_keywords": [],
        "event_keywords": [],
        "comorbidity_keywords": [],
        "special_keywords": [],
    }

    narrative = patient.get("clinical_narrative")
    is_narrative = narrative and not patient.get("primary_site")

    if is_narrative:
        _extract_from_narrative(narrative, features)
    else:
        _extract_from_structured(patient, features)

    # 去重合并
    all_kw = []
    seen = set()
    for key in features:
        for kw in features[key]:
            if kw.lower() not in seen:
                seen.add(kw.lower())
                all_kw.append(kw)
    features["all_keywords"] = all_kw

    # D2: confidence 标记
    dimensions_hit = sum(1 for k, v in features.items()
                         if k.endswith("_keywords") and k != "all_keywords" and v)
    features["confidence"] = "low" if dimensions_hit <= 2 else "high"

    return features


def _extract_from_structured(p: dict, features: dict):
    """从结构化字段提取关键词"""
    if p.get("primary_site"):
        features["diagnosis_keywords"].extend([p["primary_site"], "gastric", "胃癌"])
    if p.get("pathology"):
        features["diagnosis_keywords"].append(p["pathology"])

    for field in ("staging_prefix", "t_stage", "n_stage", "m_stage"):
        val = p.get(field)
        if val:
            features["staging_keywords"].append(val)
    prefix = p.get("staging_prefix", "")
    t = p.get("t_stage", "")
    n = p.get("n_stage", "")
    m = p.get("m_stage", "")
    if t and n:
        features["staging_keywords"].append(f"{prefix}{t}{n}{m}")

    if p.get("m_sites"):
        sites = re.split(r'[,，、/]', p["m_sites"])
        for s in sites:
            s = s.strip()
            if s:
                features["metastasis_keywords"].append(s)
                for en_name, aliases in _METASTASIS_SITES.items():
                    if any(a in s for a in aliases):
                        features["metastasis_keywords"].extend(aliases)
                        break
    if p.get("t4b_invasion"):
        features["metastasis_keywords"].append(p["t4b_invasion"])

    for field in ("biopsy_molecular", "gross_molecular"):
        val = p.get(field)
        if val:
            items = re.split(r'[,，]', val.replace("hj_", ""))
            features["molecular_keywords"].extend(i.strip() for i in items if i.strip())

    if p.get("prior_treatment"):
        matches = _TREATMENT_PATTERNS.findall(p["prior_treatment"])
        features["treatment_keywords"].extend(matches)
    if p.get("patient_type"):
        if "初治" in p["patient_type"]:
            features["treatment_keywords"].extend(["初治", "treatment-naive", "first-line"])
        elif "术前" in p["patient_type"] or "sq_" in p["patient_type"]:
            features["treatment_keywords"].extend(["术前治疗后", "post-neoadjuvant"])
    if p.get("response"):
        r = p["response"]
        if r not in ("不适用", None, ""):
            features["treatment_keywords"].append(r)

    if p.get("abnormal_markers"):
        markers = re.split(r'[,，、]', p["abnormal_markers"])
        features["marker_keywords"].extend(m.strip() for m in markers if m.strip())
    if p.get("marker_change"):
        features["marker_keywords"].append(p["marker_change"])

    if p.get("tumor_emergency") and p["tumor_emergency"] != "无":
        features["event_keywords"].append(p["tumor_emergency"])
        for ek in _EMERGENCY_KEYWORDS:
            if ek in p["tumor_emergency"]:
                features["event_keywords"].append(ek)

    if p.get("comorbidities"):
        features["comorbidity_keywords"].append(p["comorbidities"])
        matches = _COMORBIDITY_PATTERNS.findall(p["comorbidities"])
        features["comorbidity_keywords"].extend(matches)

    if p.get("siewert_type"):
        features["special_keywords"].extend([
            f"Siewert {p['siewert_type']}",
            "EGJ", "食管胃结合部",
        ])
    age = p.get("age")
    if age and isinstance(age, int) and age >= 75:
        features["special_keywords"].extend(["高龄", "elderly"])


def _extract_from_narrative(text: str, features: dict):
    """从 narrative 文本正则扫描所有维度"""
    site_patterns = ["胃", "食管", "结肠", "直肠", "gastric", "esophag", "colon", "rectal"]
    for sp in site_patterns:
        if sp in text or sp in text.lower():
            features["diagnosis_keywords"].append(sp)
    path_types = ["腺癌", "印戒", "鳞癌", "adenocarcinoma", "signet", "squamous"]
    for pt in path_types:
        if pt in text or pt in text.lower():
            features["diagnosis_keywords"].append(pt)

    features["staging_keywords"].extend(_STAGING_PATTERNS.findall(text))

    for cn_name, aliases in _METASTASIS_SITES.items():
        for a in aliases:
            if a in text:
                features["metastasis_keywords"].extend(aliases)
                break

    features["molecular_keywords"].extend(_MOLECULAR_PATTERNS.findall(text))
    features["treatment_keywords"].extend(_TREATMENT_PATTERNS.findall(text))

    marker_pats = ["CEA", "CA199", "CA724", "CA125", "AFP"]
    for mp in marker_pats:
        if mp in text.upper():
            features["marker_keywords"].append(mp)

    for ek in _EMERGENCY_KEYWORDS:
        if ek in text or ek in text.lower():
            features["event_keywords"].append(ek)

    features["comorbidity_keywords"].extend(_COMORBIDITY_PATTERNS.findall(text))

    siewert_match = re.search(r'Siewert\s*(?:type\s*)?([IⅠⅡⅢ123]+)', text, re.IGNORECASE)
    if siewert_match:
        features["special_keywords"].extend(["Siewert", "EGJ", "食管胃结合部"])
    age_match = re.search(r'(\d{2,3})\s*岁', text)
    if age_match and int(age_match.group(1)) >= 75:
        features["special_keywords"].extend(["高龄", "elderly"])


def estimate_tokens(text: str) -> int:
    """粗略估算 token 数（中英文混合按 1 token ≈ 4 字符）"""
    return len(text) // 4


def generate_batch_prompt(
    batch: list[dict],
    kb_profile: dict,
    kb_root: str,
    batch_idx: int,
    total_batches: int,
    output_file: str = "",
) -> str:
    """生成自包含的批次 prompt 文件内容。"""
    lines = []

    lines.append(f"# 批次 {batch_idx:03d}/{total_batches:03d} 检索任务\n")
    lines.append("<CONTEXT_RESET>")
    lines.append("请忽略此消息之前的所有检索结果和患者信息。")
    lines.append("以下是一个全新的、独立的批次任务，从零开始处理。")
    lines.append("不得引用或参考任何其他批次的结果。")
    lines.append('不得使用"同上"、"与前面类似"、"参考前述"等表述。')
    lines.append("</CONTEXT_RESET>\n")

    lines.append("<MANDATORY_RULES>")
    lines.append("1. 必须逐条执行以下所有 grep 命令，不得跳过任何组织")
    lines.append("2. 每位患者的每个指南组织都必须有检索结果")
    lines.append('3. 如果某指南未涉及该问题，记录: "该指南未涉及此临床问题"')
    lines.append("4. 输出必须为简体中文")
    lines.append("5. 可以补充脚本未生成的关键词，但不得删减已有的 grep 命令")
    lines.append("</MANDATORY_RULES>\n")

    lines.append(f"## 知识库\n路径: {kb_root}\n")
    lines.append("根索引:\n---")
    lines.append(kb_profile.get("root_index_content", ""))
    lines.append("---\n")

    lines.append(f"## 患者列表（本批次 {len(batch)} 人）\n")

    for pi, patient in enumerate(batch, 1):
        pid = patient.get("patient_id", "?")
        pname = patient.get("patient_name", "?")
        features = patient.get("features", {})
        grep_cmds = patient.get("grep_commands", [])

        lines.append(f"### 患者 {pi}: {pname} ({pid})\n")

        lines.append("**临床信息:**")
        for k, v in patient.items():
            if k in ("features", "grep_commands") or v is None:
                continue
            lines.append(f"- {k}: {v}")

        confidence = features.get("confidence", "high")
        lines.append(f"\n**脚本提取置信度**: {confidence}")
        if confidence == "low":
            lines.append("⚠ 该患者信息稀疏，请从临床叙述中补充推断关键词并扩展检索范围。")

        lines.append("\n**脚本提取的关键词:**")
        for dim_key in sorted(features.keys()):
            if dim_key.endswith("_keywords") and dim_key != "all_keywords":
                kws = features[dim_key]
                if kws:
                    dim_name = dim_key.replace("_keywords", "")
                    lines.append(f"- {dim_name}: {', '.join(kws)}")

        if grep_cmds:
            lines.append(f"\n#### 必须执行的 grep 命令（共 {len(grep_cmds)} 条，不得跳过）\n")
            current_org = None
            cmd_num = 0
            for gc in grep_cmds:
                if gc["org"] != current_org:
                    current_org = gc["org"]
                    lines.append(f"**{current_org}（必须）:**")
                cmd_num += 1
                lines.append(f"{cmd_num}. {gc['command']}")

        lines.append("\n#### LLM 补充检索空间")
        lines.append("如果在上述 grep 结果中发现需要进一步深入的线索，")
        lines.append("可补充执行额外的 grep 命令（限同一知识库范围内）。\n")

    lines.append("## 输出要求\n")
    lines.append(f"- 文件路径: {output_file}")
    lines.append("- 格式: JSON")
    lines.append("- 每个 guideline_results 条目必须包含: guideline, version, recommendation, evidence_level, source_file, source_lines")
    lines.append("- 输出语言: 简体中文")

    return "\n".join(lines)


def cmd_orchestrate(args):
    """orchestrate 子命令入口 — 自动编排批处理流程"""
    kb_root = resolve_kb_root(getattr(args, 'kb_root', None))
    print(f"知识库路径: {kb_root}")

    kb_profile = scan_knowledge_base(kb_root)
    if not kb_profile["orgs"]:
        print("知识库为空（无有效 org 目录）", file=sys.stderr)
        sys.exit(1)

    patients_path = Path(args.patients).resolve()
    if not patients_path.exists():
        print(f"患者文件不存在: {patients_path}", file=sys.stderr)
        sys.exit(1)
    patients_data = json.loads(patients_path.read_text(encoding="utf-8"))
    patients = patients_data.get("patients", [])
    if not patients:
        print("患者列表为空", file=sys.stderr)
        sys.exit(1)

    enriched_patients = []
    total_grep = 0
    total_kw = 0
    for p in patients:
        features = extract_patient_features(p)
        grep_cmds = generate_grep_commands(features, kb_profile, kb_root)
        enriched = {**p, "features": features, "grep_commands": grep_cmds}
        enriched_patients.append(enriched)
        total_grep += len(grep_cmds)
        total_kw += len(features.get("all_keywords", []))

    batch_size = args.batch_size
    max_tokens = args.max_prompt_tokens
    batches = _split_patients(enriched_patients, batch_size)

    final_batches = []
    for batch in batches:
        prompt = generate_batch_prompt(batch, kb_profile, str(kb_root),
                                       len(final_batches) + 1, len(batches))
        tokens = estimate_tokens(prompt)
        if tokens > max_tokens and len(batch) > 1:
            sub_batches = _auto_split_batch(batch, kb_profile, str(kb_root), max_tokens)
            final_batches.extend(sub_batches)
        else:
            final_batches.append(batch)

    output_dir = Path(args.output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    existing_plan_path = output_dir / "orchestration_plan.json"
    if existing_plan_path.exists():
        try:
            old_plan = json.loads(existing_plan_path.read_text(encoding="utf-8"))
            old_batches = old_plan.get("batches", [])
            pending = [b for b in old_batches if b.get("status") == "pending"]
            completed = [b for b in old_batches if b.get("status") == "completed"]
            if pending:
                print(f"  ℹ 检测到已有计划: {len(completed)} 批已完成, {len(pending)} 批待处理。进入续跑模式。")
            elif completed:
                print(f"  ℹ 检测到已有计划且全部完成 ({len(completed)} 批)。将重新生成。")
                print(f"     如需保留旧结果，请指定不同的 --output-dir。")
        except (json.JSONDecodeError, KeyError):
            print(f"  ⚠ 已有 orchestration_plan.json 格式异常，将重新生成。")

    plan_batches = []
    for bi, batch in enumerate(final_batches, 1):
        batch_id = f"batch_{bi:03d}"
        prompt_file = output_dir / f"{batch_id}_prompt.md"
        output_file = output_dir / f"rag_{batch_id}.json"

        status = "pending"
        if output_file.exists():
            try:
                check = json.loads(output_file.read_text(encoding="utf-8"))
                if check.get("results"):
                    status = "completed"
                    print(f"  ✓ {batch_id} 已完成 (checkpoint)")
            except (json.JSONDecodeError, KeyError):
                pass

        if status == "pending":
            prompt = generate_batch_prompt(
                batch, kb_profile, str(kb_root), bi, len(final_batches),
                output_file=str(output_file),
            )
            prompt_file.write_text(prompt, encoding="utf-8")

        plan_batches.append({
            "id": batch_id,
            "prompt_file": str(prompt_file),
            "output_file": str(output_file),
            "patients": [p.get("patient_id") for p in batch],
            "status": status,
        })

    plan = {
        "version": "2.2",
        "created_at": datetime.now().isoformat(),
        "kb_root": str(kb_root),
        "kb_profile": {
            "orgs": kb_profile["orgs"],
            "org_files": kb_profile["org_files"],
        },
        "total_patients": len(patients),
        "batch_size": batch_size,
        "batches": plan_batches,
        "next_steps": [
            f"python scripts/batch_pipeline.py merge --input-dir {output_dir} --output {output_dir.parent / 'rag_results.json'}",
            f"python scripts/batch_pipeline.py validate --input {output_dir.parent / 'rag_results.json'} --patients {patients_path}",
            f"python scripts/batch_pipeline.py generate --input {output_dir.parent / 'rag_results.json'} --format all",
        ],
        "stats": {
            "total_grep_commands": total_grep,
            "orgs_covered": sorted(kb_profile["orgs"]),
            "avg_keywords_per_patient": round(total_kw / len(patients), 1) if patients else 0,
        },
    }
    plan_path = output_dir / "orchestration_plan.json"
    plan_path.write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding="utf-8")

    pending = sum(1 for b in plan_batches if b["status"] == "pending")
    completed = sum(1 for b in plan_batches if b["status"] == "completed")
    print(f"\n编排完成:")
    print(f"  患者: {len(patients)}")
    print(f"  批次: {len(final_batches)} (待处理: {pending}, 已完成: {completed})")
    print(f"  grep 命令总数: {total_grep}")
    print(f"  组织覆盖: {', '.join(sorted(kb_profile['orgs']))}")
    print(f"  计划文件: {plan_path}")


def _auto_split_batch(
    batch: list[dict],
    kb_profile: dict,
    kb_root: str,
    max_tokens: int,
) -> list[list[dict]]:
    """D3: 递归拆分超限批次"""
    if len(batch) <= 1:
        return [batch]

    mid = len(batch) // 2
    left, right = batch[:mid], batch[mid:]
    result = []

    for sub in (left, right):
        prompt = generate_batch_prompt(sub, kb_profile, kb_root, 1, 999)
        if estimate_tokens(prompt) > max_tokens and len(sub) > 1:
            result.extend(_auto_split_batch(sub, kb_profile, kb_root, max_tokens))
        else:
            result.append(sub)

    return result


# ─── merge 子命令 ─────────────────────────────────────────────────────────────


def cmd_merge(args):
    """merge 子命令入口 — 合并多个 rag_batch_*.json 为 rag_results.json"""
    input_dir = Path(args.input_dir).resolve()
    batch_files = sorted(input_dir.glob("rag_batch_*.json"))

    if not batch_files:
        print(f"未找到批次结果文件 (rag_batch_*.json) → {input_dir}", file=sys.stderr)
        sys.exit(1)

    all_results = []
    patient_ids_seen = set()

    for bf in batch_files:
        batch_data = json.loads(bf.read_text(encoding="utf-8"))
        for result in batch_data.get("results", []):
            pid = result.get("patient_id")
            if pid in patient_ids_seen:
                print(f"  ⚠ 跳过重复患者: {pid} (来自 {bf.name})")
                continue
            patient_ids_seen.add(pid)

            # 结构规范化：确保 consensus/differences 在 clinical_questions 内
            for q in result.get("clinical_questions", []):
                if "consensus" not in q:
                    q["consensus"] = result.get("consensus", [])
                if "differences" not in q:
                    q["differences"] = result.get("differences", [])
                if "guideline_results" not in q:
                    q["guideline_results"] = result.get("guideline_results", [])
            # 清理根级别的冗余字段（规范化后不再需要）
            for key in ("consensus", "differences", "guideline_results"):
                result.pop(key, None)

            all_results.append(result)

    merged = {
        "generated_at": str(date.today()),
        "patient_count": len(all_results),
        "source_batches": [bf.name for bf in batch_files],
        "results": all_results,
    }

    output_path = Path(args.output).resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(
        json.dumps(merged, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(f"已合并 {len(batch_files)} 个批次，{len(all_results)} 位患者 → {output_path}")


# ─── validate 子命令 ──────────────────────────────────────────────────────────


def cmd_validate(args):
    """validate 子命令入口 — 检查 rag_results.json 质量与完整性"""
    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"文件不存在: {input_path}", file=sys.stderr)
        sys.exit(1)

    data = json.loads(input_path.read_text(encoding="utf-8"))
    results = data.get("results", [])

    errors = []
    warnings = []

    # 与 patients.json 对比完整性
    if args.patients:
        patients_path = Path(args.patients).resolve()
        if patients_path.exists():
            patients_data = json.loads(patients_path.read_text(encoding="utf-8"))
            expected_ids = {p["patient_id"] for p in patients_data.get("patients", [])}
            actual_ids = {r.get("patient_id") for r in results}
            missing = expected_ids - actual_ids
            extra = actual_ids - expected_ids
            if missing:
                errors.append(f"缺失患者 ({len(missing)}): {', '.join(sorted(missing))}")
            if extra:
                warnings.append(f"多余患者 ({len(extra)}): {', '.join(sorted(extra))}")

    # 逐患者检查
    rec_lengths = []
    for r in results:
        pid = r.get("patient_id", "?")

        # 必要字段
        for field in ("diagnosis_summary", "clinical_questions", "disease_type"):
            if not r.get(field):
                errors.append(f"[{pid}] 缺失字段: {field}")

        questions = r.get("clinical_questions", [])
        if not questions:
            errors.append(f"[{pid}] 无临床问题")
            rec_lengths.append((pid, 0))
            continue

        total_len = 0
        for qi, q in enumerate(questions, 1):
            grs = q.get("guideline_results", [])
            if not grs:
                warnings.append(f"[{pid}] Q{qi} 无指南检索结果")

            for g in grs:
                rec = g.get("recommendation", "")
                total_len += len(rec)
                if len(rec) < 50:
                    warnings.append(
                        f"[{pid}] Q{qi} {g.get('guideline', '')} 推荐过短 ({len(rec)}字)"
                    )
                if not g.get("evidence_level"):
                    warnings.append(
                        f"[{pid}] Q{qi} {g.get('guideline', '')} 缺失证据等级"
                    )
                if not g.get("source_file"):
                    warnings.append(
                        f"[{pid}] Q{qi} {g.get('guideline', '')} 缺失来源文件"
                    )

            if not q.get("consensus"):
                warnings.append(f"[{pid}] Q{qi} 缺失共识分析")
            if not q.get("differences"):
                warnings.append(f"[{pid}] Q{qi} 缺失差异分析")

        rec_lengths.append((pid, total_len))

    # 跨患者一致性：检测质量下降
    if len(rec_lengths) >= 3:
        lengths = [l for _, l in rec_lengths if l > 0]
        if lengths:
            avg_len = sum(lengths) / len(lengths)
            for pid, length in rec_lengths:
                if length > 0 and length < avg_len * 0.3:
                    warnings.append(
                        f"[{pid}] 推荐总长度异常偏短 ({length}字 vs 平均 {avg_len:.0f}字)"
                    )

    # 输出报告
    print(f"验证结果: {len(results)} 位患者")
    if errors:
        print(f"\n  ✗ {len(errors)} 个错误:")
        for e in errors:
            print(f"    ✗ {e}")
    if warnings:
        print(f"\n  ⚠ {len(warnings)} 个警告:")
        for w in warnings:
            print(f"    ⚠ {w}")
    if not errors and not warnings:
        print(f"  ✓ 验证通过，数据完整")

    sys.exit(1 if errors else 0)


# ─── generate 子命令 ──────────────────────────────────────────────────────────


def load_rag_results(path: Path) -> dict:
    """加载 RAG 结果 JSON"""
    text = path.read_text(encoding="utf-8")
    return json.loads(text)


def generate_xlsx(data: dict, output_path: Path):
    """生成批量推荐汇总表"""
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = "批量推荐汇总"

    # 收集所有指南名称（动态列）
    all_guidelines = []
    for result in data.get("results", []):
        for q in result.get("clinical_questions", []):
            for gr in q.get("guideline_results", []):
                name = gr.get("guideline", "")
                if name and name not in all_guidelines:
                    all_guidelines.append(name)

    # 表头
    fixed_headers = ["患者ID", "姓名", "肿瘤部位", "诊断摘要", "临床问题"]
    summary_headers = ["共识点", "差异点"]
    guideline_headers = [f"{g}推荐" for g in all_guidelines]
    tail_headers = ["备注"]
    headers = fixed_headers + summary_headers + guideline_headers + tail_headers

    # 样式
    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    alt_fill = PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid")
    wrap_align = Alignment(wrap_text=True, vertical="top")

    # 写表头
    for col_idx, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = wrap_align

    # 写数据行
    row_num = 2
    for result in data.get("results", []):
        questions = result.get("clinical_questions", [])
        question_text = "; ".join(q.get("question", "") for q in questions)

        # 按指南名汇总推荐
        guideline_recs = {}
        consensus = []
        differences = []
        for q in questions:
            for gr in q.get("guideline_results", []):
                g_name = gr.get("guideline", "")
                rec = gr.get("recommendation", "")
                level = gr.get("evidence_level", "")
                entry = f"{rec} ({level})" if level else rec
                guideline_recs.setdefault(g_name, []).append(entry)
            consensus.extend(q.get("consensus", []))
            differences.extend(q.get("differences", []))

        row_data = [
            result.get("patient_id", ""),
            result.get("patient_name", ""),
            result.get("primary_site", ""),
            result.get("diagnosis_summary", ""),
            question_text,
            "\n".join(consensus) if consensus else "—",
            "\n".join(differences) if differences else "—",
        ]
        for g in all_guidelines:
            recs = guideline_recs.get(g, [])
            row_data.append("\n".join(recs) if recs else "—")
        row_data.append("")  # 备注

        for col_idx, val in enumerate(row_data, 1):
            cell = ws.cell(row=row_num, column=col_idx, value=val)
            cell.alignment = wrap_align
            if row_num % 2 == 0:
                cell.fill = alt_fill

        row_num += 1

    # 列宽
    col_widths = {"患者ID": 16, "姓名": 10, "肿瘤部位": 14, "诊断摘要": 35,
                  "临床问题": 40, "共识点": 35, "差异点": 35, "备注": 20}
    for col_idx, header in enumerate(headers, 1):
        width = col_widths.get(header, 35)
        ws.column_dimensions[get_column_letter(col_idx)].width = width

    # 冻结首行 + 自动筛选
    ws.freeze_panes = "A2"
    ws.auto_filter.ref = ws.dimensions

    output_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(output_path)
    print(f"  ✓ 汇总表: {output_path}")


def generate_docx(data: dict, output_dir: Path):
    """生成个体推荐意见书 (每人一份 DOCX)"""
    from docx import Document
    from docx.enum.section import WD_ORIENT
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor

    output_dir.mkdir(parents=True, exist_ok=True)
    count = 0

    for result in data.get("results", []):
        doc = Document()

        # 页面设置 — Landscape
        section = doc.sections[0]
        section.orientation = WD_ORIENT.LANDSCAPE
        new_width, new_height = section.page_height, section.page_width
        section.page_width = new_width
        section.page_height = new_height
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(0.8)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)

        pid = result.get("patient_id", "未知")
        pname = result.get("patient_name", "未知")

        # 标题
        title = doc.add_heading(f"{pname}（{pid}）临床指南推荐意见书", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # 1. 患者信息
        doc.add_heading("1. 患者信息", level=1)
        info_fields = [
            ("患者ID", pid),
            ("姓名", pname),
            ("肿瘤部位", result.get("primary_site", "—")),
            ("病种诊断", result.get("disease_type", "—")),
            ("诊断摘要", result.get("diagnosis_summary", "—")),
        ]
        info_table = doc.add_table(rows=len(info_fields), cols=2)
        info_table.style = "Table Grid"
        info_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        for i, (label, value) in enumerate(info_fields):
            info_table.cell(i, 0).text = label
            info_table.cell(i, 1).text = str(value) if value else "—"
            # 加粗标签列
            for run in info_table.cell(i, 0).paragraphs[0].runs:
                run.bold = True

        # 2-N. 每个临床问题
        for qi, q in enumerate(result.get("clinical_questions", []), 1):
            doc.add_heading(f"2.{qi} 临床问题", level=1)
            doc.add_paragraph(q.get("question", ""))

            # 推荐对比表
            doc.add_heading(f"2.{qi}.1 各指南推荐对比", level=2)
            guideline_results = q.get("guideline_results", [])
            if guideline_results:
                tbl = doc.add_table(rows=1 + len(guideline_results), cols=5)
                tbl.style = "Table Grid"
                tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
                # 表头
                for ci, header in enumerate(["指南", "版本", "推荐意见", "证据等级", "来源"]):
                    cell = tbl.cell(0, ci)
                    cell.text = header
                    for run in cell.paragraphs[0].runs:
                        run.bold = True
                # 数据行
                for ri, gr in enumerate(guideline_results, 1):
                    tbl.cell(ri, 0).text = gr.get("guideline", "")
                    tbl.cell(ri, 1).text = gr.get("version", "")
                    tbl.cell(ri, 2).text = gr.get("recommendation", "")
                    tbl.cell(ri, 3).text = gr.get("evidence_level", "")
                    source = gr.get("source_file", "")
                    lines = gr.get("source_lines", "")
                    tbl.cell(ri, 4).text = f"{source} 第{lines}行" if lines else source

            # 共识与差异
            doc.add_heading(f"2.{qi}.2 指南间共识与差异", level=2)
            consensus = q.get("consensus", [])
            differences = q.get("differences", [])
            if consensus:
                doc.add_paragraph("共识点:", style="List Bullet")
                for c in consensus:
                    doc.add_paragraph(c, style="List Bullet 2")
            if differences:
                doc.add_paragraph("主要差异:", style="List Bullet")
                for d in differences:
                    doc.add_paragraph(d, style="List Bullet 2")

        # 生成日期
        date_para = doc.add_paragraph(f"生成日期: {data.get('generated_at', str(date.today()))}")
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in date_para.runs:
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(128, 128, 128)

        # 免责声明
        doc.add_paragraph("")
        disclaimer = doc.add_paragraph(
            "本文档由医学指南RAG系统自动生成，仅供临床参考，不替代专业医学判断。"
        )
        disclaimer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in disclaimer.runs:
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(128, 128, 128)

        # 保存
        safe_name = pname.replace("/", "_").replace("\\", "_")
        filename = f"{pid}_{safe_name}_推荐意见书.docx"
        doc.save(output_dir / filename)
        count += 1

    print(f"  ✓ 推荐意见书: {count} 份 → {output_dir}/")


def generate_pptx(data: dict, output_path: Path):
    """生成批量推荐幻灯片（使用模板）"""
    import math

    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    template_path = Path(__file__).parent.parent / "templates" / "report_template.pptx"
    prs = Presentation(str(template_path))

    # ─── 删除模板中的占位 slide ───
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    results = data.get("results", [])

    # ─── Helper: 截断文本 ───
    def _truncate(text, max_chars=150):
        if not text or len(text) <= max_chars:
            return text
        # 在 max_chars 范围内找最近的句号/分号/逗号
        for sep in ["。", "；", "，", ".", ";", ","]:
            pos = text.rfind(sep, 0, max_chars)
            if pos > max_chars // 2:
                return text[:pos + 1] + "…"
        return text[:max_chars] + "…"

    # ─── Helper: 设置表格单元格样式 ───
    def _style_cell(cell, font_size, bold=False, bg_color=None, font_color=None):
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = bold
            if font_color:
                p.font.color.rgb = font_color
        cell.text_frame.word_wrap = True
        if bg_color:
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = parse_xml(f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:srgbClr val="{bg_color}"/></a:solidFill>')
            tcPr.append(solidFill)

    # ─── Helper: 向 textbox 添加带格式的段落 ───
    def _add_para(tf, text, font_size, bold=False, color=None, first=False):
        if first:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        if color:
            p.font.color.rgb = color
        return p

    # ─── Helper: 给 textbox 设置白色不透明背景 ───
    def _fill_white(shape):
        sp = shape._element
        spPr = sp.find(qn('a:spPr'))
        if spPr is None:
            spPr = parse_xml(f'<a:spPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
            sp.append(spPr)
        solidFill = parse_xml(
            '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:srgbClr val="FFFFFF"/></a:solidFill>'
        )
        spPr.insert(0, solidFill)

    # ─── Slide 1: 封面 (Layout 0 — Title Slide) ───
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[0].text = "批量指南推荐报告"
    for run in slide.placeholders[0].text_frame.paragraphs[0].runs:
        run.font.size = Pt(44)
        run.font.bold = True
    subtitle_tf = slide.placeholders[1].text_frame
    subtitle_tf.paragraphs[0].text = f"日期: {data.get('generated_at', str(date.today()))}  |  患者数: {len(results)}"
    for run in subtitle_tf.paragraphs[0].runs:
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # ─── Slide 2+: 汇总概览 (Layout 1 — Custom Layout) ───
    max_rows_per_page = 15
    total_pages = math.ceil(len(results) / max_rows_per_page) if results else 1

    for page_idx in range(total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if total_pages > 1:
            slide.placeholders[0].text = f"患者汇总概览（{page_idx + 1}/{total_pages}）"
        else:
            slide.placeholders[0].text = "患者汇总概览"

        page_results = results[page_idx * max_rows_per_page : (page_idx + 1) * max_rows_per_page]
        rows_count = len(page_results) + 1  # +1 for header
        cols_count = 4
        tbl_shape = slide.shapes.add_table(
            rows_count, cols_count,
            Inches(0.92), Inches(0.85), Inches(11.50), Inches(5.95)
        )
        tbl = tbl_shape.table

        # 设置列宽: ID 1.8in, 姓名 1.2in, 诊断摘要 5.0in, 核心推荐 3.5in
        col_widths = [Inches(1.8), Inches(1.2), Inches(5.0), Inches(3.5)]
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = w

        # 表头
        for ci, header in enumerate(["ID", "姓名", "诊断摘要", "核心推荐"]):
            cell = tbl.cell(0, ci)
            cell.text = header
            _style_cell(cell, 10, bold=True, bg_color="1F4E79",
                        font_color=RGBColor(0xFF, 0xFF, 0xFF))

        # 数据行
        for ri, result in enumerate(page_results, 1):
            tbl.cell(ri, 0).text = result.get("patient_id", "")
            tbl.cell(ri, 1).text = result.get("patient_name", "")
            tbl.cell(ri, 2).text = _truncate(result.get("diagnosis_summary", ""), 60)
            questions = result.get("clinical_questions", [])
            first_rec = ""
            if questions:
                grs = questions[0].get("guideline_results", [])
                if grs:
                    first_rec = _truncate(grs[0].get("recommendation", ""), 60)
            tbl.cell(ri, 3).text = first_rec

            bg = "F2F2F2" if ri % 2 == 0 else None
            for ci in range(cols_count):
                _style_cell(tbl.cell(ri, ci), 9, bg_color=bg,
                            font_color=RGBColor(0x33, 0x33, 0x33))

    # ─── 每患者 3 页 ───
    # Layout 背景元素位置参考:
    #   Layout 2: 5×2 info table (0.92,1.27,11.50,2.30) + "临床问题" textbox (0.92,4.76,11.50,1.01)
    #   Layout 3: "各指南推荐对比" textbox (0.92,0.97,11.50,0.40) + 4×4 table (0.92,1.51,11.50,2.05)
    #   Layout 4: 无非占位符背景 → 适合共识差异页

    for result in results:
        pid = result.get("patient_id", "")
        pname = result.get("patient_name", "")
        site = result.get("primary_site", "—")
        disease = result.get("disease_type", "—")
        summary = result.get("diagnosis_summary", "—")
        questions = result.get("clinical_questions", [])

        # ── Slide A: 患者信息 + 临床问题 (Layout 2) ──
        # Layout 2 有背景 5×2 表格和 "临床问题" 文本框，需在 slide 上添加覆盖层
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        slide.placeholders[0].text = f"患者 {pid}  {pname}"

        # 覆盖 Layout 的 5×2 info table: 精确匹配位置 (0.92,1.27,11.50,2.30)
        info_tbl_shape = slide.shapes.add_table(
            5, 2, Inches(0.92), Inches(1.27), Inches(11.50), Inches(2.30)
        )
        info_tbl = info_tbl_shape.table
        info_tbl.columns[0].width = Inches(1.80)
        info_tbl.columns[1].width = Inches(9.70)
        info_data = [
            ("患者ID", pid),
            ("姓名", pname),
            ("肿瘤部位", site),
            ("病种诊断", disease),
            ("诊断摘要", summary),
        ]
        for ri, (label, value) in enumerate(info_data):
            info_tbl.cell(ri, 0).text = label
            info_tbl.cell(ri, 1).text = str(value) if value else "—"
            _style_cell(info_tbl.cell(ri, 0), 13, bold=True,
                        bg_color="1F4E79", font_color=RGBColor(0xFF, 0xFF, 0xFF))
            _style_cell(info_tbl.cell(ri, 1), 13,
                        bg_color="F2F2F2" if ri % 2 == 0 else "FFFFFF",
                        font_color=RGBColor(0x33, 0x33, 0x33))

        # Layout 自带 "临床问题：" 标签 textbox 在 (0.92,4.76,11.50,1.01)
        # 在其下方填入实际的临床问题内容
        q_box = slide.shapes.add_textbox(
            Inches(0.92), Inches(5.15), Inches(11.50), Inches(1.65)
        )
        q_tf = q_box.text_frame
        q_tf.word_wrap = True
        q_tf.auto_size = MSO_AUTO_SIZE.NONE
        info_color = RGBColor(0x33, 0x33, 0x33)
        for qi, q in enumerate(questions, 1):
            _add_para(q_tf, f"{qi}. {q.get('question', '')}", 12, color=info_color, first=(qi == 1))

        # ── Slide B: 推荐对比表 (Layout 3) ──
        # Layout 3 有背景文字和 4×4 表格，用更大的表格覆盖
        slide = prs.slides.add_slide(prs.slide_layouts[3])
        slide.placeholders[0].text = f"各指南推荐对比 — {pname}"

        all_grs = []
        for q in questions:
            all_grs.extend(q.get("guideline_results", []))

        if all_grs:
            data_font = 9 if len(all_grs) <= 8 else 8
            header_font = 10 if len(all_grs) <= 8 else 9

            rows_n = 1 + len(all_grs)
            # 覆盖 Layout 的文字和表格: 从 y=0.90 开始，覆盖 textbox(0.97) 和 table(1.51)
            tbl_shape = slide.shapes.add_table(
                rows_n, 4,
                Inches(0.92), Inches(0.90), Inches(11.50), Inches(5.90)
            )
            tbl = tbl_shape.table
            for ci, w in enumerate([Inches(1.3), Inches(1.0), Inches(7.2), Inches(2.0)]):
                tbl.columns[ci].width = w

            for ci, header in enumerate(["指南", "版本", "推荐意见", "证据等级"]):
                cell = tbl.cell(0, ci)
                cell.text = header
                _style_cell(cell, header_font, bold=True, bg_color="1F4E79",
                            font_color=RGBColor(0xFF, 0xFF, 0xFF))

            for ri, gr in enumerate(all_grs, 1):
                tbl.cell(ri, 0).text = gr.get("guideline", "")
                tbl.cell(ri, 1).text = gr.get("version", "")
                tbl.cell(ri, 2).text = _truncate(gr.get("recommendation", ""), 150)
                tbl.cell(ri, 3).text = gr.get("evidence_level", "")
                bg = "F2F2F2" if ri % 2 == 0 else "FFFFFF"
                for ci in range(4):
                    _style_cell(tbl.cell(ri, ci), data_font, bg_color=bg,
                                font_color=RGBColor(0x33, 0x33, 0x33))

        # ── Slide C: 共识与差异 (Layout 4 — Comparison，无背景装饰) ──
        slide = prs.slides.add_slide(prs.slide_layouts[4])
        slide.placeholders[0].text = f"指南间共识与差异 — {pname}"

        consensus_all = []
        differences_all = []
        for q in questions:
            consensus_all.extend(q.get("consensus", []))
            differences_all.extend(q.get("differences", []))

        total_items = len(consensus_all) + len(differences_all) + 2
        body_font = 11 if total_items <= 18 else 10
        title_color = RGBColor(0x1F, 0x4E, 0x79)
        body_color = RGBColor(0x33, 0x33, 0x33)

        # 左列：共识点 — ph[1] 标签保持, ph[2] 填入内容
        slide.placeholders[1].text = "共识点"
        left_tf = slide.placeholders[2].text_frame
        left_tf.word_wrap = True
        left_tf.auto_size = MSO_AUTO_SIZE.NONE
        left_tf.clear()
        for ci, c in enumerate(consensus_all):
            _add_para(left_tf, f"• {c}", body_font, color=body_color, first=(ci == 0))

        # 右列：主要差异 — ph[3] 标签保持, ph[4] 填入内容
        slide.placeholders[3].text = "主要差异"
        right_tf = slide.placeholders[4].text_frame
        right_tf.word_wrap = True
        right_tf.auto_size = MSO_AUTO_SIZE.NONE
        right_tf.clear()
        for di, d in enumerate(differences_all):
            _add_para(right_tf, f"• {d}", body_font, color=body_color, first=(di == 0))

    # ─── 最后一页: 免责声明 (Layout 0 — Title Slide，无背景装饰) ───
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[0].text = "免责声明"
    subtitle_tf = slide.placeholders[1].text_frame
    subtitle_tf.paragraphs[0].text = "本幻灯片由医学指南RAG系统自动生成，仅供临床参考，不替代专业医学判断。"
    for run in subtitle_tf.paragraphs[0].runs:
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    total_slides = 1 + total_pages + 3 * len(results) + 1
    print(f"  ✓ 幻灯片: {output_path} ({total_slides} 页)")


def cmd_generate(args):
    """generate 子命令入口"""
    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"RAG 结果文件不存在: {input_path}", file=sys.stderr)
        sys.exit(1)

    data = load_rag_results(input_path)
    patient_count = data.get("patient_count", len(data.get("results", [])))
    print(f"加载 RAG 结果: {patient_count} 位患者\n")

    output_dir = Path(args.output_dir).resolve()
    fmt = args.format

    if fmt in ("all", "xlsx"):
        generate_xlsx(data, output_dir / "批量推荐汇总表.xlsx")
    if fmt in ("all", "docx"):
        generate_docx(data, output_dir / "reports")
    if fmt in ("all", "pptx"):
        generate_pptx(data, output_dir / "批量推荐幻灯片.pptx")

    print(f"\n生成完成 → {output_dir}/")


# ─── CLI ──────────────────────────────────────────────────────────────────────


def main():
    parser = argparse.ArgumentParser(
        description="批量患者指南检索管道工具",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    sub = parser.add_subparsers(dest="command", required=True)

    # parse
    p_parse = sub.add_parser("parse", help="解析输入 xlsx → patients.json")
    p_parse.add_argument("--input", required=True, help="输入 xlsx 文件路径")
    p_parse.add_argument("--output", default="Output/patients.json", help="输出 JSON 路径")

    # split
    p_split = sub.add_parser("split", help="将 patients.json 分成多个批次文件")
    p_split.add_argument("--input", required=True, help="patients.json 路径")
    p_split.add_argument("--batch-size", type=int, default=5, help="每批患者数 (默认 5)")
    p_split.add_argument("--output-dir", default="Output/batches", help="批次文件输出目录")

    # orchestrate
    p_orch = sub.add_parser("orchestrate", help="自动编排批处理流程（扫描知识库+生成 prompt）")
    p_orch.add_argument("--patients", required=True, help="patients.json 路径")
    p_orch.add_argument("--kb-root", default=None, help="知识库根路径（可选）")
    p_orch.add_argument("--output-dir", default="Output/batches", help="输出目录")
    p_orch.add_argument("--batch-size", type=int, default=5, help="每批患者数 (默认 5)")
    p_orch.add_argument("--max-prompt-tokens", type=int, default=80000,
                         help="单个 prompt 最大 token 数 (默认 80000)")

    # merge
    p_merge = sub.add_parser("merge", help="合并批次结果为 rag_results.json")
    p_merge.add_argument("--input-dir", required=True, help="批次结果所在目录")
    p_merge.add_argument("--output", default="Output/rag_results.json", help="合并输出路径")

    # validate
    p_validate = sub.add_parser("validate", help="验证 RAG 结果质量与完整性")
    p_validate.add_argument("--input", required=True, help="rag_results.json 路径")
    p_validate.add_argument("--patients", help="patients.json 路径（可选，用于完整性对比）")

    # generate
    p_gen = sub.add_parser("generate", help="从 RAG 结果生成产出物")
    p_gen.add_argument("--input", required=True, help="RAG 结果 JSON 路径")
    p_gen.add_argument("--output-dir", default="Output", help="输出目录")
    p_gen.add_argument("--format", choices=["all", "xlsx", "docx", "pptx"], default="all",
                       help="输出格式 (默认 all)")

    args = parser.parse_args()
    if args.command == "parse":
        cmd_parse(args)
    elif args.command == "split":
        cmd_split(args)
    elif args.command == "orchestrate":
        cmd_orchestrate(args)
    elif args.command == "merge":
        cmd_merge(args)
    elif args.command == "validate":
        cmd_validate(args)
    elif args.command == "generate":
        cmd_generate(args)


if __name__ == "__main__":
    main()
